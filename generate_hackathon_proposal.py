#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
2025 新竹政策黑客松奪冠提案簡報生成器
赤土崎全齡社福樞紐 - Championship-Winning Hackathon Proposal
10 Slides, 16:9, Professional Layout with Open Source Images
"""

import sys
import io

# 修復 Windows 控制台編碼問題
if sys.platform == 'win32':
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
    sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import requests
from io import BytesIO

# ========== 色彩主題設計 ==========
PRIMARY = RGBColor(0, 255, 136)      # 霓虹綠 #00ff88 (創新科技)
SECONDARY = RGBColor(0, 204, 255)    # 青藍色 #00ccff (專業可信)
ACCENT = RGBColor(255, 170, 0)       # 橙色 #ffaa00 (熱情活力)
WARNING = RGBColor(255, 100, 100)    # 紅色 (緊迫性)
SUCCESS = RGBColor(100, 255, 100)    # 亮綠 (成功)
BG_DARK = RGBColor(10, 14, 39)       # 深藍背景 #0a0e27
BG_CARD = RGBColor(20, 25, 50)       # 卡片背景
TEXT_WHITE = RGBColor(255, 255, 255)
TEXT_GRAY = RGBColor(180, 180, 180)

# ========== 開源圖片 URL ==========
IMAGES = {
    'cover': 'https://images.unsplash.com/photo-1486406146926-c627a92ad1ab?w=1920&q=80',  # 現代建築
    'problem': 'https://images.unsplash.com/photo-1551836022-deb4988cc6c0?w=1920&q=80',  # 壓力工作
    'scattered': 'https://images.unsplash.com/photo-1524661135-423995f22d0b?w=1920&q=80',  # 城市交通
    'building': 'https://images.unsplash.com/photo-1545324418-cc1a3fa10c00?w=1920&q=80',  # 綠色建築
    'innovation': 'https://images.unsplash.com/photo-1451187580459-43490279c0fa?w=1920&q=80',  # 科技創新
    'democracy': 'https://images.unsplash.com/photo-1529107386315-e1a2ed48a620?w=1920&q=80',  # 社區參與
    'finance': 'https://images.unsplash.com/photo-1460925895917-afdab827c52f?w=1920&q=80',  # 財務成長
    'family': 'https://images.unsplash.com/photo-1609220136736-443140cffec6?w=1920&q=80',  # 三代同堂
    'global': 'https://images.unsplash.com/photo-1526778548025-fa2f459cd5c1?w=1920&q=80',  # 地球全球
    'vision': 'https://images.unsplash.com/photo-1511895426328-dc8714191300?w=1920&q=80',  # 未來願景
}

def download_image(url, max_retries=2):
    """從URL下載圖片，有重試機制"""
    for attempt in range(max_retries):
        try:
            response = requests.get(url, timeout=15)
            if response.status_code == 200:
                return BytesIO(response.content)
            else:
                print(f"  ⚠️  圖片下載失敗 (狀態: {response.status_code}): {url[:60]}...")
        except Exception as e:
            print(f"  ❌ 圖片下載錯誤 (嘗試 {attempt+1}/{max_retries}): {e}")
    return None

def add_background_dark(slide):
    """添加深色漸層背景"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK

def add_image_with_overlay(slide, image_url, overlay_opacity=0.4):
    """添加背景圖片並加上半透明遮罩"""
    img_stream = download_image(image_url)
    if img_stream:
        try:
            # 添加圖片作為背景
            pic = slide.shapes.add_picture(
                img_stream,
                Inches(0), Inches(0),
                width=Inches(10),
                height=Inches(5.625)
            )
            # 移到最底層
            slide.shapes._spTree.remove(pic._element)
            slide.shapes._spTree.insert(2, pic._element)

            # 添加半透明黑色遮罩
            overlay = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0),
                Inches(10), Inches(5.625)
            )
            overlay.fill.solid()
            overlay.fill.fore_color.rgb = BG_DARK
            overlay.fill.transparency = overlay_opacity
            overlay.line.fill.background()

            return True
        except Exception as e:
            print(f"  ❌ 背景圖片設置失敗: {e}")
    return False

def add_title_box(slide, text, top, color=PRIMARY, size=44, bold=True):
    """添加標題文字框"""
    box = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(9), Inches(0.8))
    tf = box.text_frame
    tf.text = text
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.LEFT
    return box

def add_stat_box(slide, left, top, width, label, value, color=PRIMARY):
    """添加統計數據框"""
    # 背景卡片
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top),
        Inches(width), Inches(1.2)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = BG_CARD
    card.line.color.rgb = color
    card.line.width = Pt(2)

    # 數值
    value_box = slide.shapes.add_textbox(
        Inches(left), Inches(top + 0.2),
        Inches(width), Inches(0.5)
    )
    tf = value_box.text_frame
    tf.text = value
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER

    # 標籤
    label_box = slide.shapes.add_textbox(
        Inches(left), Inches(top + 0.75),
        Inches(width), Inches(0.3)
    )
    tf = label_box.text_frame
    tf.text = label
    p = tf.paragraphs[0]
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_GRAY
    p.alignment = PP_ALIGN.CENTER

# ========== 初始化簡報 (16:9) ==========
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

print("🎨 開始生成 2025 新竹政策黑客松奪冠提案簡報...")
print(f"   尺寸: 16:9 (10\" × 5.625\")")
print(f"   總頁數: 10 張投影片\n")

# ==================== 第 1 頁：封面 ====================
print("📄 [1/10] 生成封面頁...")
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
add_image_with_overlay(slide1, IMAGES['cover'], 0.5)

# 主標題
title = slide1.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
tf = title.text_frame
tf.text = "赤土崎全齡社福樞紐"
p = tf.paragraphs[0]
p.font.size = Pt(66)
p.font.bold = True
p.font.color.rgb = PRIMARY
p.alignment = PP_ALIGN.CENTER

# 副標題
subtitle = slide1.shapes.add_textbox(Inches(0.5), Inches(2.7), Inches(9), Inches(0.6))
tf = subtitle.text_frame
tf.text = "竹科家庭全齡支持中心 - 解決時間貧窮危機"
p = tf.paragraphs[0]
p.font.size = Pt(28)
p.font.color.rgb = SECONDARY
p.alignment = PP_ALIGN.CENTER

# 標籤
tag = slide1.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(9), Inches(0.5))
tf = tag.text_frame
tf.text = "2025 新竹政策黑客松提案 | 標案 114A109 | 設計預算 2,287 萬"
p = tf.paragraphs[0]
p.font.size = Pt(16)
p.font.color.rgb = TEXT_GRAY
p.alignment = PP_ALIGN.CENTER

# ==================== 第 2 頁：問題陳述 ====================
print("📄 [2/10] 生成問題陳述頁...")
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_image_with_overlay(slide2, IMAGES['problem'], 0.6)

# 標題
add_title_box(slide2, "💔 年薪 300 萬的「甲級貧戶」", 0.3, WARNING, 48)

# 內容區域
content = slide2.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(8.6), Inches(3.5))
tf = content.text_frame
tf.word_wrap = True

items = [
    ("時間貧窮", "每日接送 2 小時 × 250 天 = 500 小時/年", "相當於 20 天家庭時間被通勤吃掉"),
    ("服務分散", "長照在東區、托嬰在北區、課輔在香山", "3 個地點 × 3 趟往返 = 精疲力盡"),
    ("托育困難", "公托排隊 16-20 個月，私托 3-5 萬/月", "雙薪家庭也負擔沉重"),
    ("心理壓力", "竹科社群焦慮、育兒孤立、婚姻緊張", "36.4% 文獻提及關係危機"),
]

for i, (title, stat, detail) in enumerate(items):
    p = tf.add_paragraph()
    p.text = f"🔸 {title}"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.space_before = Pt(16) if i > 0 else Pt(0)

    p = tf.add_paragraph()
    p.text = f"   • {stat}"
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_WHITE
    p.level = 1

    p = tf.add_paragraph()
    p.text = f"   • {detail}"
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_GRAY
    p.level = 1

# ==================== 第 3 頁：根本原因 ====================
print("📄 [3/10] 生成根本原因頁...")
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_background_dark(slide3)

add_title_box(slide3, "🎯 根本原因：服務分散 ≠ 缺服務", 0.3, PRIMARY, 40)

# 核心洞見
insight = slide3.shapes.add_textbox(Inches(1), Inches(1.2), Inches(8), Inches(0.8))
tf = insight.text_frame
tf.text = "「竹科家庭的痛點不是『缺服務』，而是『服務分散』」"
p = tf.paragraphs[0]
p.font.size = Pt(24)
p.font.italic = True
p.font.color.rgb = SECONDARY
p.alignment = PP_ALIGN.CENTER

# 對比表格
comparison_data = [
    ("現況模式", "整合方案", "改善幅度"),
    ("3 個地點", "1 棟建築", "↓ 67% 通勤"),
    ("年補助 5,500萬", "年補助 0 萬 (Year 3)", "↓ 100% 依賴"),
    ("各自管理", "垂直整合", "↑ 300% 空間效率"),
    ("2 小時/日", "30 分鐘/日", "↑ 1.5 小時家庭時間"),
]

top = 2.3
for i, (col1, col2, col3) in enumerate(comparison_data):
    is_header = (i == 0)
    bg_color = ACCENT if is_header else BG_CARD
    text_color = BG_DARK if is_header else TEXT_WHITE
    font_size = 16 if is_header else 14

    for j, text in enumerate([col1, col2, col3]):
        box = slide3.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.7 + j * 3), Inches(top + i * 0.55),
            Inches(2.8), Inches(0.5)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = bg_color
        box.line.fill.background()

        text_box = slide3.shapes.add_textbox(
            Inches(0.7 + j * 3), Inches(top + i * 0.55),
            Inches(2.8), Inches(0.5)
        )
        tf = text_box.text_frame
        tf.text = text
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.font.size = Pt(font_size)
        p.font.bold = is_header
        p.font.color.rgb = text_color
        p.alignment = PP_ALIGN.CENTER

# ==================== 第 4 頁：解決方案 ====================
print("📄 [4/10] 生成解決方案頁...")
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_image_with_overlay(slide4, IMAGES['building'], 0.6)

add_title_box(slide4, "🏢 解決方案：垂直整合 One-Stop 服務", 0.3, SUCCESS, 38)

# 樓層配置
floors = [
    ("B1", "停車場 + 機電", "30 車位 (含無障礙)", BG_CARD),
    ("1F", "長照日照中心", "50-60 位長輩/日", PRIMARY),
    ("2F", "公共托嬰中心", "40-50 位嬰幼兒/日", SECONDARY),
    ("3F", "家庭支持服務", "諮商 + 社區廚房 + 志工", ACCENT),
    ("4F", "青少年活動中心", "30-40 位青少年/日", RGBColor(150, 100, 255)),
]

left_start = 1.5
top_start = 1.4
for i, (floor, name, detail, color) in enumerate(floors):
    # 樓層標籤
    floor_box = slide4.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left_start), Inches(top_start + i * 0.7),
        Inches(0.8), Inches(0.6)
    )
    floor_box.fill.solid()
    floor_box.fill.fore_color.rgb = color
    floor_box.line.fill.background()

    text_box = slide4.shapes.add_textbox(
        Inches(left_start), Inches(top_start + i * 0.7),
        Inches(0.8), Inches(0.6)
    )
    tf = text_box.text_frame
    tf.text = floor
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.alignment = PP_ALIGN.CENTER

    # 功能名稱
    name_box = slide4.shapes.add_textbox(
        Inches(left_start + 1), Inches(top_start + i * 0.7),
        Inches(3.5), Inches(0.3)
    )
    tf = name_box.text_frame
    tf.text = name
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE

    # 詳細說明
    detail_box = slide4.shapes.add_textbox(
        Inches(left_start + 1), Inches(top_start + i * 0.7 + 0.32),
        Inches(3.5), Inches(0.25)
    )
    tf = detail_box.text_frame
    tf.text = detail
    p = tf.paragraphs[0]
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_GRAY

# 右側統計
add_stat_box(slide4, 6.5, 1.5, 2.8, "總樓地板面積", "3,100 m²", SUCCESS)
add_stat_box(slide4, 6.5, 2.9, 2.8, "每日服務人次", "140-180 人", PRIMARY)

# ==================== 第 5 頁：6 大技術創新 ====================
print("📄 [5/10] 生成技術創新頁...")
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_image_with_overlay(slide5, IMAGES['innovation'], 0.65)

add_title_box(slide5, "⚡ 6 大技術創新（全台首創）", 0.3, PRIMARY, 40)

innovations = [
    ("1. 垂直隔音設計", "4F 籃球場 ≠ 干擾 1F 失智專區 (50dB 隔音)", PRIMARY),
    ("2. 時間錯峰服務", "120 座位餐廳 → 150 人/日使用 (125% 效率)", SECONDARY),
    ("3. AI 記憶重生", "失智長輩生命故事 AI 重現 (MMSE +2.8 分)", ACCENT),
    ("4. 預測性健康", "72% 準確率預測 7 日內生病風險", WARNING),
    ("5. SIDS 監測系統", "98% 靈敏度嬰兒猝死預防 (50 床即時監控)", RGBColor(255, 100, 255)),
    ("6. 時間銀行自動化", "志工配對 2 天 → 15 分鐘 (92% 成功率)", SUCCESS),
]

top = 1.3
for i, (title, desc, color) in enumerate(innovations):
    row = i // 2
    col = i % 2

    card = slide5.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5 + col * 4.8), Inches(top + row * 1.15),
        Inches(4.5), Inches(1)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = BG_CARD
    card.fill.transparency = 0.1
    card.line.color.rgb = color
    card.line.width = Pt(3)

    title_box = slide5.shapes.add_textbox(
        Inches(0.7 + col * 4.8), Inches(top + 0.15 + row * 1.15),
        Inches(4.1), Inches(0.35)
    )
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = color

    desc_box = slide5.shapes.add_textbox(
        Inches(0.7 + col * 4.8), Inches(top + 0.52 + row * 1.15),
        Inches(4.1), Inches(0.4)
    )
    tf = desc_box.text_frame
    tf.text = desc
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_WHITE

# ==================== 第 6 頁：政策創新 ====================
print("📄 [6/10] 生成政策創新頁...")
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_image_with_overlay(slide6, IMAGES['democracy'], 0.6)

add_title_box(slide6, "🤝 政策創新：從「政府供給」到「社區共創」", 0.3, ACCENT, 34)

policies = [
    {
        'icon': '🏢',
        'title': 'ESG 企業整合',
        'points': [
            '台積電 + 聯發科投入 2.85 億',
            '員工福利 ROI 回報',
            '企業永續報告書亮點',
        ],
        'color': PRIMARY
    },
    {
        'icon': '⏰',
        'title': '時間銀行制度',
        'points': [
            '年交易量 3,000-8,000 筆',
            '等值服務 120 萬元',
            '參與率 68% (勝過現金)',
        ],
        'color': SECONDARY
    },
    {
        'icon': '🗳️',
        'title': '參與式預算',
        'points': [
            '25% 營運預算 (250 萬)',
            '社區居民投票決定',
            '每季大會民主治理',
        ],
        'color': ACCENT
    },
]

for i, policy in enumerate(policies):
    card = slide6.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5 + i * 3.2), Inches(1.4),
        Inches(3), Inches(3.5)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = BG_CARD
    card.fill.transparency = 0.2
    card.line.color.rgb = policy['color']
    card.line.width = Pt(2)

    # 圖示
    icon_box = slide6.shapes.add_textbox(
        Inches(0.5 + i * 3.2), Inches(1.6),
        Inches(3), Inches(0.5)
    )
    tf = icon_box.text_frame
    tf.text = policy['icon']
    p = tf.paragraphs[0]
    p.font.size = Pt(48)
    p.alignment = PP_ALIGN.CENTER

    # 標題
    title_box = slide6.shapes.add_textbox(
        Inches(0.7 + i * 3.2), Inches(2.3),
        Inches(2.6), Inches(0.4)
    )
    tf = title_box.text_frame
    tf.text = policy['title']
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = policy['color']
    p.alignment = PP_ALIGN.CENTER

    # 內容
    content_box = slide6.shapes.add_textbox(
        Inches(0.7 + i * 3.2), Inches(2.8),
        Inches(2.6), Inches(2)
    )
    tf = content_box.text_frame
    tf.word_wrap = True
    for point in policy['points']:
        p = tf.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_WHITE
        p.space_before = Pt(6)

# ==================== 第 7 頁：財務模式 ====================
print("📄 [7/10] 生成財務模式頁...")
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_image_with_overlay(slide7, IMAGES['finance'], 0.65)

add_title_box(slide7, "💰 財務模式：3 年自給自足", 0.3, SUCCESS, 44)

# 5年財務預測表格
years = ["項目", "Y1 (2029)", "Y2 (2030)", "Y3 (2031)", "Y5 (2033)"]
cost = ["營運成本", "2,328萬", "2,200萬", "1,940萬", "2,010萬"]
revenue = ["自籌收入", "1,850萬", "2,750萬", "3,025萬", "3,200萬"]
subsidy = ["政府補助", "975萬", "150萬", "0萬", "0萬"]
surplus = ["盈餘", "+497萬", "+700萬", "+1,085萬", "+1,190萬"]

data = [years, cost, revenue, subsidy, surplus]

top = 1.4
for i, row in enumerate(data):
    for j, cell in enumerate(row):
        is_header_row = (i == 0)
        is_header_col = (j == 0)

        if is_header_row or is_header_col:
            bg = ACCENT
            text_color = BG_DARK
            bold = True
        else:
            bg = BG_CARD
            text_color = TEXT_WHITE
            bold = False

        box = slide7.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5 + j * 1.9), Inches(top + i * 0.5),
            Inches(1.9), Inches(0.5)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = bg
        box.line.color.rgb = ACCENT
        box.line.width = Pt(1)

        text_box = slide7.shapes.add_textbox(
            Inches(0.5 + j * 1.9), Inches(top + i * 0.5),
            Inches(1.9), Inches(0.5)
        )
        tf = text_box.text_frame
        tf.text = cell
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.font.size = Pt(11)
        p.font.bold = bold
        p.font.color.rgb = text_color
        p.alignment = PP_ALIGN.CENTER

# 底部重點
highlight = slide7.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(0.8))
tf = highlight.text_frame
tf.text = "🎯 政府節省：15 年共省 8.11 億元 (vs 傳統分散模式)\n💡 自給率：Year 3 達 156%，Year 5 達 159%"
for p in tf.paragraphs:
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = SUCCESS
    p.alignment = PP_ALIGN.CENTER

# ==================== 第 8 頁：跨代影響 ====================
print("📄 [8/10] 生成跨代影響頁...")
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
add_image_with_overlay(slide8, IMAGES['family'], 0.55)

add_title_box(slide8, "👨‍👩‍👧‍👦 跨代影響：老中青三贏", 0.3, PRIMARY, 44)

impacts = [
    {
        'group': '1F 長輩 (65+)',
        'color': PRIMARY,
        'metrics': [
            ('MMSE 認知', '+2.8 分', '(延緩失智)'),
            ('焦慮降低', '-40%', '(心理健康)'),
            ('社交孤立', '-51%', '(參與活動)'),
        ]
    },
    {
        'group': '2F 幼兒 (0-2 歲)',
        'color': SECONDARY,
        'metrics': [
            ('語言發展', '+18%', '(跨代刺激)'),
            ('安全保障', '98%', '(SIDS 監測)'),
            ('社區歸屬', '8.2/10', '(溫暖環境)'),
        ]
    },
    {
        'group': '4F 青少年 (13-18)',
        'color': ACCENT,
        'metrics': [
            ('職涯探索', 'STEM', '(創客空間)'),
            ('領導能力', '志工', '(服務學習)'),
            ('同理心', '+35%', '(長輩互動)'),
        ]
    },
]

for i, impact in enumerate(impacts):
    # 卡片
    card = slide8.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5 + i * 3.2), Inches(1.4),
        Inches(3), Inches(3.2)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = BG_CARD
    card.fill.transparency = 0.2
    card.line.color.rgb = impact['color']
    card.line.width = Pt(3)

    # 群組標題
    title_box = slide8.shapes.add_textbox(
        Inches(0.7 + i * 3.2), Inches(1.6),
        Inches(2.6), Inches(0.4)
    )
    tf = title_box.text_frame
    tf.text = impact['group']
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = impact['color']
    p.alignment = PP_ALIGN.CENTER

    # 指標
    for j, (label, value, note) in enumerate(impact['metrics']):
        # 標籤
        label_box = slide8.shapes.add_textbox(
            Inches(0.8 + i * 3.2), Inches(2.2 + j * 0.85),
            Inches(2.4), Inches(0.25)
        )
        tf = label_box.text_frame
        tf.text = label
        p = tf.paragraphs[0]
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_GRAY

        # 數值
        value_box = slide8.shapes.add_textbox(
            Inches(0.8 + i * 3.2), Inches(2.45 + j * 0.85),
            Inches(2.4), Inches(0.35)
        )
        tf = value_box.text_frame
        tf.text = value
        p = tf.paragraphs[0]
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = impact['color']
        p.alignment = PP_ALIGN.CENTER

        # 註釋
        note_box = slide8.shapes.add_textbox(
            Inches(0.8 + i * 3.2), Inches(2.78 + j * 0.85),
            Inches(2.4), Inches(0.2)
        )
        tf = note_box.text_frame
        tf.text = note
        p = tf.paragraphs[0]
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_GRAY
        p.alignment = PP_ALIGN.CENTER

# ==================== 第 9 頁：時機關鍵 ====================
print("📄 [9/10] 生成時機關鍵頁...")
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
add_background_dark(slide9)

add_title_box(slide9, "⏰ 黃金窗口：現在就是最佳時機！", 0.3, WARNING, 40)

# 時間軸
timeline = [
    ("2025/10/30", "標案開標", "設計廠商確定", WARNING),
    ("2025/11-12", "初步設計", "功能 100% 可調整 ✅", SUCCESS),
    ("2026/01-03", "細部設計", "調整難度 ⚠️ (10%)", ACCENT),
    ("2026/05+", "施工階段", "改變成本 +50-100% ❌", RGBColor(100, 100, 100)),
]

top = 1.5
for i, (date, phase, detail, color) in enumerate(timeline):
    # 時間點圓點
    circle = slide9.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(1), Inches(top + i * 0.9),
        Inches(0.4), Inches(0.4)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()

    # 日期
    date_box = slide9.shapes.add_textbox(
        Inches(1.6), Inches(top + i * 0.9),
        Inches(2), Inches(0.4)
    )
    tf = date_box.text_frame
    tf.text = date
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = color

    # 階段
    phase_box = slide9.shapes.add_textbox(
        Inches(3.8), Inches(top + i * 0.9),
        Inches(2.5), Inches(0.4)
    )
    tf = phase_box.text_frame
    tf.text = phase
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE

    # 詳情
    detail_box = slide9.shapes.add_textbox(
        Inches(6.5), Inches(top + i * 0.9),
        Inches(3), Inches(0.4)
    )
    tf = detail_box.text_frame
    tf.text = detail
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_GRAY

# 強調框
emphasis = slide9.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(1), Inches(4.8),
    Inches(8), Inches(0.6)
)
emphasis.fill.solid()
emphasis.fill.fore_color.rgb = WARNING
emphasis.line.fill.background()

text_box = slide9.shapes.add_textbox(Inches(1), Inches(4.8), Inches(8), Inches(0.6))
tf = text_box.text_frame
tf.text = "🔥 贏得黑客松 = 影響政策設計 = 改變 2.3 億建築未來"
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = BG_DARK
p.alignment = PP_ALIGN.CENTER

# ==================== 第 10 頁：願景與行動呼籲 ====================
print("📄 [10/10] 生成願景呼籲頁...")
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
add_image_with_overlay(slide10, IMAGES['vision'], 0.5)

# 願景標題
vision_title = slide10.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(0.7))
tf = vision_title.text_frame
tf.text = "🌟 從「政府供給」到「社區共創」"
p = tf.paragraphs[0]
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = PRIMARY
p.alignment = PP_ALIGN.CENTER

# 承諾
promises = [
    "👵 長輩找回生命故事 (AI 記憶重生)",
    "👨‍👩‍👧 父母找回家庭時間 (30 分 vs 2 小時)",
    "👶 孩子找回社區溫暖 (跨代共學)",
    "🏛️ 政府找回財政責任 (Year 3 自給)",
]

promise_box = slide10.shapes.add_textbox(Inches(1.5), Inches(2), Inches(7), Inches(2))
tf = promise_box.text_frame
tf.word_wrap = True
for i, promise in enumerate(promises):
    p = tf.add_paragraph()
    p.text = promise
    p.font.size = Pt(20)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(12) if i > 0 else Pt(0)
    p.alignment = PP_ALIGN.LEFT

# 致勝公式
formula = slide10.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(0.5))
tf = formula.text_frame
tf.text = "高門檻 + 政策對接 + 普世共鳴 + 可落地 = 🏆 奪冠"
p = tf.paragraphs[0]
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 215, 0)  # 金色
p.alignment = PP_ALIGN.CENTER

# CTA
cta = slide10.shapes.add_textbox(Inches(1), Inches(4.8), Inches(8), Inches(0.6))
tf = cta.text_frame
tf.text = "選擇赤土崎 — 讓台灣引領亞洲社福轉型"
p = tf.paragraphs[0]
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = SECONDARY
p.alignment = PP_ALIGN.CENTER

# ========== 儲存簡報 ==========
output = "2025新竹政策黑客松_赤土崎全齡社福樞紐_奪冠提案.pptx"
prs.save(output)

print(f"\n✅ 簡報生成成功！")
print(f"📁 檔案: {output}")
print(f"📐 尺寸: 16:9 (標準簡報比例)")
print(f"📄 總頁數: 10 張投影片")
print(f"\n🎯 簡報結構:")
print("   1. 封面 - 赤土崎全齡社福樞紐")
print("   2. 問題 - 年薪 300 萬的甲級貧戶")
print("   3. 原因 - 服務分散 ≠ 缺服務")
print("   4. 方案 - 垂直整合 One-Stop")
print("   5. 創新 - 6 大技術突破")
print("   6. 政策 - 社區共創模式")
print("   7. 財務 - 3 年自給自足")
print("   8. 影響 - 老中青三贏")
print("   9. 時機 - 黃金窗口")
print("   10. 願景 - 台灣引領亞洲")
print(f"\n💡 圖片來源: Unsplash (免費商用)")
print(f"🎨 設計風格: 深色科技風 + 霓虹配色")
print(f"🏆 目標: 2025 新竹政策黑客松冠軍\n")
