#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
赤土崎多功能館 - 8大前瞻創新設計簡報生成器
Generate Professional PowerPoint Presentation with Images from Open Source Libraries
"""

import sys
import io

# 修復 Windows 控制台編碼問題
if sys.platform == 'win32':
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
    sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import requests
from io import BytesIO
from PIL import Image

# 初始化簡報
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# 定義色彩主題（赤土崎品牌色）
PRIMARY_COLOR = RGBColor(0, 255, 136)  # 霓虹綠 #00ff88
SECONDARY_COLOR = RGBColor(0, 204, 255)  # 青藍色 #00ccff
ACCENT_COLOR = RGBColor(255, 0, 255)  # 洋紅色 #ff00ff
BG_DARK = RGBColor(10, 14, 39)  # 深藍背景 #0a0e27
TEXT_COLOR = RGBColor(255, 255, 255)  # 白色文字

# 開源圖片URL列表（來自 Unsplash 和 Pexels）
IMAGE_URLS = {
    'cover': 'https://images.unsplash.com/photo-1486406146926-c627a92ad1ab?w=1200',  # 現代建築
    'ai_elderly': 'https://images.unsplash.com/photo-1581579438747-1dc8d17bbce4?w=1200',  # AI + 長照
    'time_bank': 'https://images.unsplash.com/photo-1559027615-cd4628902d4a?w=1200',  # 社區志工
    'sids_monitor': 'https://images.unsplash.com/photo-1515488042361-ee00e0ddd4e4?w=1200',  # 嬰兒照護
    'intergenerational': 'https://images.unsplash.com/photo-1581579186913-45ac3e6efe93?w=1200',  # 老少互動
    'vr_education': 'https://images.unsplash.com/photo-1622979135225-d2ba269cf1ac?w=1200',  # VR 教育
    'stem_lab': 'https://images.unsplash.com/photo-1581092160607-ee22621dd758?w=1200',  # STEM 實驗室
    'solar_energy': 'https://images.pexels.com/photos/356036/pexels-photo-356036.jpeg?w=1200',  # 太陽能板
    'esg_exhibition': 'https://images.unsplash.com/photo-1557804506-669a67965ba0?w=1200',  # ESG 展示
}

def download_image(url):
    """從 URL 下載圖片"""
    try:
        response = requests.get(url, timeout=10)
        if response.status_code == 200:
            return BytesIO(response.content)
        else:
            print(f"⚠️  無法下載圖片: {url} (狀態碼: {response.status_code})")
            return None
    except Exception as e:
        print(f"❌ 圖片下載失敗: {url} - {e}")
        return None

def add_background(slide, color=BG_DARK):
    """添加深色背景"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide(prs):
    """創建封面頁"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版面
    add_background(slide)

    # 標題
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "赤土崎多功能館"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(60)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_COLOR
    title_para.alignment = PP_ALIGN.CENTER

    # 副標題
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "8大前瞻創新設計"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(36)
    subtitle_para.font.color.rgb = SECONDARY_COLOR
    subtitle_para.alignment = PP_ALIGN.CENTER

    # 下方資訊
    info_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1))
    info_frame = info_box.text_frame
    info_frame.text = "新竹市114年政策黑客松提案\n竹科家庭全齡支持中心"
    for para in info_frame.paragraphs:
        para.font.size = Pt(20)
        para.font.color.rgb = TEXT_COLOR
        para.alignment = PP_ALIGN.CENTER

    # 嘗試添加封面圖片（作為背景）
    img_stream = download_image(IMAGE_URLS['cover'])
    if img_stream:
        try:
            left = Inches(0)
            top = Inches(0)
            pic = slide.shapes.add_picture(img_stream, left, top, width=prs.slide_width)
            # 移到最底層
            slide.shapes._spTree.remove(pic._element)
            slide.shapes._spTree.insert(2, pic._element)

            # 添加半透明遮罩
            shape = slide.shapes.add_shape(1, left, top, prs.slide_width, prs.slide_height)
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(10, 14, 39)
            shape.fill.transparency = 0.3
        except Exception as e:
            print(f"添加封面圖片失敗: {e}")

def add_content_slide(prs, title, content_dict, image_url=None, theme_color=PRIMARY_COLOR):
    """創建內容頁"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)

    # 標題區塊
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = theme_color
    title_para.alignment = PP_ALIGN.LEFT

    # 圖片區域（左側）
    img_left = Inches(0.5)
    img_top = Inches(1.3)
    img_width = Inches(4.5)
    img_height = Inches(5.7)

    # 嘗試下載並添加圖片
    img_added = False
    if image_url:
        img_stream = download_image(image_url)
        if img_stream:
            try:
                slide.shapes.add_picture(img_stream, img_left, img_top, width=img_width, height=img_height)
                img_added = True
            except Exception as e:
                print(f"添加圖片失敗 ({title}): {e}")

    # 如果圖片未添加，創建彩色佔位符
    if not img_added:
        placeholder = slide.shapes.add_shape(1, img_left, img_top, img_width, img_height)
        fill = placeholder.fill
        fill.solid()
        fill.fore_color.rgb = theme_color

        # 添加圖片說明文字
        text_box = slide.shapes.add_textbox(img_left, img_top + img_height/2 - Inches(0.3), img_width, Inches(0.6))
        tf = text_box.text_frame
        tf.text = "📷 請從 Unsplash/Pexels\n免費圖庫替換此圖片"
        for para in tf.paragraphs:
            para.font.size = Pt(14)
            para.font.color.rgb = RGBColor(255, 255, 255)
            para.alignment = PP_ALIGN.CENTER

    # 內容區域（右側）
    content_left = Inches(5.2)
    content_top = Inches(1.3)
    content_width = Inches(4.3)

    # 創建內容文字框
    content_box = slide.shapes.add_textbox(content_left, content_top, content_width, Inches(5.7))
    tf = content_box.text_frame
    tf.word_wrap = True

    # 添加各項內容
    for key, value in content_dict.items():
        p = tf.add_paragraph()
        p.text = f"{key}"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = theme_color
        p.space_before = Pt(12)

        if isinstance(value, list):
            for item in value:
                p = tf.add_paragraph()
                p.text = f"• {item}"
                p.font.size = Pt(12)
                p.font.color.rgb = TEXT_COLOR
                p.level = 1
        else:
            p = tf.add_paragraph()
            p.text = str(value)
            p.font.size = Pt(12)
            p.font.color.rgb = TEXT_COLOR

def add_summary_slide(prs):
    """創建總結頁"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)

    # 標題
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "🎯 致勝公式"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_COLOR
    title_para.alignment = PP_ALIGN.CENTER

    # 公式內容
    formula_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
    tf = formula_box.text_frame
    tf.word_wrap = True

    formula_items = [
        "高門檻 (社福專業 + 空間規劃)",
        "+",
        "政策對接 (標案 2,287萬 + 財劃法 214億)",
        "+",
        "普世共鳴 (竹科家庭痛點)",
        "+",
        "可落地性 (8大創新技術實證)",
        "=",
        "🏆 新竹政策黑客松冠軍"
    ]

    for i, item in enumerate(formula_items):
        p = tf.add_paragraph()
        p.text = item
        if item in ['+', '=']:
            p.font.size = Pt(32)
            p.font.color.rgb = SECONDARY_COLOR
            p.alignment = PP_ALIGN.CENTER
        elif '🏆' in item:
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 215, 0)  # 金色
            p.alignment = PP_ALIGN.CENTER
        else:
            p.font.size = Pt(20)
            p.font.color.rgb = TEXT_COLOR
            p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(8)

# ========== 生成簡報內容 ==========

print("🎨 開始生成赤土崎多功能館簡報...")

# 1. 封面頁
print("  📄 生成封面頁...")
add_title_slide(prs)

# 2. 創新設計內容頁
innovations = [
    {
        'title': '1️⃣ AI記憶重生系統',
        'content': {
            '樓層位置': '1F 失智症專區 (156 m²)',
            '核心技術': ['AI感測設備', '互動投影系統', '記憶喚醒算法'],
            '創新亮點': '全台首創AI輔助失智症記憶療法，結合感測與投影技術協助長輩喚醒珍貴記憶',
            '設備規格': ['感測設備×12', '互動投影×4', '活動室', '無障礙衛浴'],
            '政策對接': '長照2.0 - 失智症照護服務',
        },
        'image': 'ai_elderly',
        'color': PRIMARY_COLOR
    },
    {
        'title': '2️⃣ 時間銀行',
        'content': {
            '樓層位置': '1F 共享空間',
            '核心概念': '志工時數儲存與兌換，建立社區互助經濟體系',
            '創新亮點': ['突破傳統志工模式', '創造社區貨幣', '永續互助機制'],
            '設施配置': ['時數兌換台', '檔案櫃系統', '交流討論區', '數位管理平台'],
            '參與式預算': '5F大會堂可容納250人進行社區治理討論',
        },
        'image': 'time_bank',
        'color': RGBColor(255, 170, 0)
    },
    {
        'title': '3️⃣ SIDS嬰兒猝死監測系統',
        'content': {
            '樓層位置': '2F 嬰兒教室 (120 m²)',
            '技術規格': '98%靈敏度的即時生理監測系統',
            '創新亮點': '🏆 新竹市首座配備完整SIDS系統的公共托嬰中心',
            '監測設備': ['SIDS監測器×50組', '嬰兒床×50張', '中央監控系統', '緊急警報'],
            '監測項目': ['呼吸頻率', '心跳', '體溫', '血氧濃度', '睡姿'],
        },
        'image': 'sids_monitor',
        'color': RGBColor(255, 100, 100)
    },
    {
        'title': '4️⃣ 跨代互動室',
        'content': {
            '樓層位置': '3F 中央區域 (49.5 m²)',
            '核心價值': '打破世代隔閡，創造社會資本',
            '活動類型': ['藝術共創', '技藝傳承', '生命故事', '節慶慶典'],
            '設施配置': ['工作台×3', '座位×12', '藝術用品', '展示區'],
            '實證效益': '研究顯示跨代互動可提升長輩生活品質，並增強青少年同理心與自信',
        },
        'image': 'intergenerational',
        'color': RGBColor(255, 0, 255)
    },
    {
        'title': '5️⃣ VR共學教室',
        'content': {
            '樓層位置': '4F 投影展示區 (36 m²)',
            '核心設備': ['VR頭盔×10', '4K攝影棚', '綠幕系統', '編輯工作站×4'],
            '創新亮點': '沉浸式虛擬實境學習 + 影音創作編輯一體化',
            '應用場景': ['歷史文化體驗', '科學實驗模擬', 'Podcast錄製', '影音剪輯教學'],
            '技術優勢': 'VR學習記憶保留率比傳統教學提升9%',
        },
        'image': 'vr_education',
        'color': SECONDARY_COLOR
    },
    {
        'title': '6️⃣ STEM科技教室',
        'content': {
            '樓層位置': '4F 北西區 (45 m²)',
            '核心設備': ['Arduino×10', 'Raspberry Pi×10', '3D列印機×2', '機器人套件×10'],
            '創新亮點': '完整創客教育生態系統，培養新世代工程師',
            '教學模組': ['程式設計', '電子電路', '3D建模列印', '機器人競賽'],
            '產業連結': '串聯竹科企業資源，提供實習與職涯探索機會',
        },
        'image': 'stem_lab',
        'color': RGBColor(0, 255, 200)
    },
    {
        'title': '7️⃣ 太陽能綠能系統',
        'content': {
            '樓層位置': '7F 頂層 (99 m²)',
            '系統規格': ['太陽能板×6片', '發電設備', '即時監控系統', '儲能系統'],
            '創新亮點': '建築能源自主，響應淨零碳排目標',
            '環境效益': ['年發電量約3.5萬度', '減碳約18公噸/年', '節省電費約10萬/年'],
            '教育價值': '結合STEM教室進行綠能教育，培養環境永續意識',
        },
        'image': 'solar_energy',
        'color': RGBColor(100, 200, 0)
    },
    {
        'title': '8️⃣ ESG企業展示廳',
        'content': {
            '樓層位置': '6F (240 m²)',
            '核心功能': '串聯竹科企業CSR資源，建立公私協力生態系',
            '展示內容': ['企業永續實踐', 'DEI多元共融', '社會創新專案', '員工志工成果'],
            '設施配置': ['展示櫃×8', '互動螢幕', '業師辦公室', '會議室'],
            '商業模式': '企業贊助換取展示空間，資源投入社福服務',
        },
        'image': 'esg_exhibition',
        'color': RGBColor(150, 100, 255)
    },
]

for i, innovation in enumerate(innovations):
    print(f"  📄 生成第 {i+2} 頁: {innovation['title']}")
    add_content_slide(
        prs,
        innovation['title'],
        innovation['content'],
        IMAGE_URLS.get(innovation['image']),
        innovation['color']
    )

# 3. 總結頁
print("  📄 生成總結頁...")
add_summary_slide(prs)

# 保存簡報
output_file = '赤土崎多功能館_8大前瞻創新設計.pptx'
prs.save(output_file)
print(f"\n✅ 簡報生成成功！")
print(f"📁 檔案位置: {output_file}")
print(f"\n📸 圖片來源說明:")
print("  • Unsplash: https://unsplash.com (免費商用，無需署名)")
print("  • Pexels: https://www.pexels.com (免費商用，無需署名)")
print("\n💡 建議: 請自行從上述網站下載高品質圖片替換簡報中的佔位符圖片")
