#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
聯發科 ESG 合作提案簡報生成器 V2.0 (16:9 重新設計版)
生成日期: 2025-10-23
設計原則: 視覺優先、一頁一概念、數據視覺化
報告人: 蔡秀吉
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime

def create_mediatek_presentation_v2():
    """創建聯發科 ESG 合作提案簡報 V2.0 (16:9)"""

    # 創建簡報對象 - 16:9 比例
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 標準寬度
    prs.slide_height = Inches(7.5)

    # 定義顏色方案（聯發科品牌色）
    MTK_ORANGE = RGBColor(255, 102, 0)
    NAVY_BLUE = RGBColor(0, 51, 102)
    LIGHT_BLUE = RGBColor(100, 181, 246)
    WHITE = RGBColor(255, 255, 255)
    LIGHT_GRAY = RGBColor(245, 245, 245)
    DARK_GRAY = RGBColor(66, 66, 66)
    SUCCESS_GREEN = RGBColor(76, 175, 80)
    WARNING_ORANGE = RGBColor(255, 152, 0)

    # 定義留白與網格
    MARGIN_LEFT = Inches(1.0)
    MARGIN_RIGHT = Inches(1.0)
    MARGIN_TOP = Inches(0.75)
    MARGIN_BOTTOM = Inches(0.75)
    CONTENT_WIDTH = Inches(11.333)
    CONTENT_HEIGHT = Inches(6.0)

    # ========== Slide 1: 封面 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 背景漸層（深藍到淺藍）
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = NAVY_BLUE

    # 頂部紅色裝飾條
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.333), Inches(0.15)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = MTK_ORANGE
    top_bar.line.fill.background()

    # 主標題
    title_box = slide.shapes.add_textbox(
        MARGIN_LEFT, Inches(2.5), CONTENT_WIDTH, Inches(1.2)
    )
    title_frame = title_box.text_frame
    title_frame.text = "竹科家庭照護時數銀行"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(66)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.alignment = PP_ALIGN.CENTER

    # 副標題（核心價值主張）
    subtitle_box = slide.shapes.add_textbox(
        MARGIN_LEFT, Inches(4.0), CONTENT_WIDTH, Inches(0.8)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "投資 500 萬，創造 1,020 萬社會價值"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(36)
    subtitle_para.font.color.rgb = LIGHT_BLUE
    subtitle_para.alignment = PP_ALIGN.CENTER

    # 合作方
    partner_box = slide.shapes.add_textbox(
        MARGIN_LEFT, Inches(5.2), CONTENT_WIDTH, Inches(0.5)
    )
    partner_frame = partner_box.text_frame
    partner_frame.text = "聯發科文教基金會 × 新竹市政府"
    partner_para = partner_frame.paragraphs[0]
    partner_para.font.size = Pt(24)
    partner_para.font.color.rgb = WHITE
    partner_para.alignment = PP_ALIGN.CENTER

    # 日期
    date_box = slide.shapes.add_textbox(
        MARGIN_LEFT, Inches(6.5), CONTENT_WIDTH, Inches(0.4)
    )
    date_frame = date_box.text_frame
    date_frame.text = "2025 年 10 月"
    date_para = date_frame.paragraphs[0]
    date_para.font.size = Pt(18)
    date_para.font.color.rgb = LIGHT_GRAY
    date_para.alignment = PP_ALIGN.CENTER

    # ========== Slide 2: 議程 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 標題
    title_box = slide.shapes.add_textbox(
        MARGIN_LEFT, MARGIN_TOP, CONTENT_WIDTH, Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = "簡報大綱"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = NAVY_BLUE

    # 議程卡片
    agenda_items = [
        ("01", "聯發科面臨的挑戰", "人才留任 × 招募困難 × ESG落後"),
        ("02", "赤土崎解決方案", "一站式整合服務模式"),
        ("03", "投資效益分析", "SROI 1:2.04 社會價值"),
        ("04", "合作方案", "3年投資計畫與回報"),
        ("05", "立即行動", "11/15前決策窗口")
    ]

    y_start = Inches(1.8)
    card_height = Inches(0.85)
    card_spacing = Inches(0.1)

    for i, (num, title, desc) in enumerate(agenda_items):
        y_pos = y_start + i * (card_height + card_spacing)

        # 卡片背景
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            MARGIN_LEFT, y_pos, CONTENT_WIDTH, card_height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_GRAY if i % 2 == 0 else WHITE
        card.line.color.rgb = NAVY_BLUE
        card.line.width = Pt(1)

        # 編號圓圈
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            MARGIN_LEFT + Inches(0.3), y_pos + Inches(0.18), Inches(0.5), Inches(0.5)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = MTK_ORANGE
        circle.line.fill.background()
        circle_tf = circle.text_frame
        circle_tf.text = num
        circle_tf.paragraphs[0].font.size = Pt(20)
        circle_tf.paragraphs[0].font.bold = True
        circle_tf.paragraphs[0].font.color.rgb = WHITE
        circle_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        circle_tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        # 標題
        title_tf = slide.shapes.add_textbox(
            MARGIN_LEFT + Inches(1.0), y_pos + Inches(0.15),
            Inches(6), Inches(0.35)
        )
        title_tf.text_frame.text = title
        title_tf.text_frame.paragraphs[0].font.size = Pt(24)
        title_tf.text_frame.paragraphs[0].font.bold = True
        title_tf.text_frame.paragraphs[0].font.color.rgb = NAVY_BLUE

        # 描述
        desc_tf = slide.shapes.add_textbox(
            MARGIN_LEFT + Inches(1.0), y_pos + Inches(0.5),
            Inches(10), Inches(0.3)
        )
        desc_tf.text_frame.text = desc
        desc_tf.text_frame.paragraphs[0].font.size = Pt(16)
        desc_tf.text_frame.paragraphs[0].font.color.rgb = DARK_GRAY

    # ========== Slide 3: 聯發科的挑戰（概覽）==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 標題
    title_box = slide.shapes.add_textbox(
        MARGIN_LEFT, MARGIN_TOP, CONTENT_WIDTH, Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = "聯發科面臨的三大人才挑戰"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = NAVY_BLUE

    # 三大挑戰卡片（左中右排列）
    challenges = [
        ("💰", "員工家庭壓力", "54.5%", "時間貧窮", "離職成本\n30萬/人"),
        ("👥", "招募吸引力不足", "-10%", "Offer接受率", "人才外流\n竹科 vs 矽谷"),
        ("📊", "ESG評級落後", "BBB", "MSCI評級", "融資成本↑\n產業平均A")
    ]

    card_width = Inches(3.5)
    card_height = Inches(4.0)
    card_spacing = Inches(0.4)
    x_start = MARGIN_LEFT + Inches(0.2)
    y_start = Inches(2.0)

    for i, (emoji, title, metric, metric_desc, impact) in enumerate(challenges):
        x_pos = x_start + i * (card_width + card_spacing)

        # 卡片背景
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x_pos, y_start, card_width, card_height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_GRAY
        card.line.color.rgb = MTK_ORANGE
        card.line.width = Pt(2)

        # Emoji/Icon
        emoji_tf = slide.shapes.add_textbox(
            x_pos + Inches(0.2), y_start + Inches(0.2),
            card_width - Inches(0.4), Inches(0.6)
        )
        emoji_tf.text_frame.text = emoji
        emoji_tf.text_frame.paragraphs[0].font.size = Pt(48)
        emoji_tf.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # 標題
        title_tf = slide.shapes.add_textbox(
            x_pos + Inches(0.2), y_start + Inches(0.9),
            card_width - Inches(0.4), Inches(0.6)
        )
        title_tf.text_frame.text = title
        title_tf.text_frame.paragraphs[0].font.size = Pt(20)
        title_tf.text_frame.paragraphs[0].font.bold = True
        title_tf.text_frame.paragraphs[0].font.color.rgb = NAVY_BLUE
        title_tf.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # 關鍵指標（大數字）
        metric_tf = slide.shapes.add_textbox(
            x_pos + Inches(0.2), y_start + Inches(1.7),
            card_width - Inches(0.4), Inches(0.8)
        )
        metric_tf.text_frame.text = metric
        metric_tf.text_frame.paragraphs[0].font.size = Pt(52)
        metric_tf.text_frame.paragraphs[0].font.bold = True
        metric_tf.text_frame.paragraphs[0].font.color.rgb = MTK_ORANGE
        metric_tf.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # 指標描述
        desc_tf = slide.shapes.add_textbox(
            x_pos + Inches(0.2), y_start + Inches(2.5),
            card_width - Inches(0.4), Inches(0.4)
        )
        desc_tf.text_frame.text = metric_desc
        desc_tf.text_frame.paragraphs[0].font.size = Pt(16)
        desc_tf.text_frame.paragraphs[0].font.color.rgb = DARK_GRAY
        desc_tf.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # 影響說明
        impact_tf = slide.shapes.add_textbox(
            x_pos + Inches(0.2), y_start + Inches(3.1),
            card_width - Inches(0.4), Inches(0.7)
        )
        impact_tf.text_frame.text = impact
        impact_tf.text_frame.paragraphs[0].font.size = Pt(16)
        impact_tf.text_frame.paragraphs[0].font.color.rgb = DARK_GRAY
        impact_tf.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        impact_tf.text_frame.word_wrap = True

    # 底部總結
    summary_box = slide.shapes.add_textbox(
        MARGIN_LEFT, Inches(6.3), CONTENT_WIDTH, Inches(0.5)
    )
    summary_frame = summary_box.text_frame
    summary_frame.text = "➜  年度隱形成本：離職損失 1.2 億 + 招募成本 8,000 萬 + ESG 融資成本增加"
    summary_para = summary_frame.paragraphs[0]
    summary_para.font.size = Pt(20)
    summary_para.font.bold = True
    summary_para.font.color.rgb = MTK_ORANGE
    summary_para.alignment = PP_ALIGN.CENTER

    # ========== Slide 4: 挑戰1深入 - 員工家庭壓力 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 標題
    title_box = slide.shapes.add_textbox(
        MARGIN_LEFT, MARGIN_TOP, CONTENT_WIDTH, Inches(0.7)
    )
    title_frame = title_box.text_frame
    title_frame.text = "挑戰 1：員工家庭照顧壓力 → 離職率上升"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = NAVY_BLUE

    # 左側：三明治世代圖示（文字版流程圖）
    left_box = slide.shapes.add_textbox(
        MARGIN_LEFT, Inches(1.8), Inches(5.5), Inches(4.2)
    )
    left_frame = left_box.text_frame
    left_frame.text = """👴 長輩照顧（65+）
需要日照、失智照護

        ↕

👨‍💼 竹科工程師（35歲）
「三明治世代」壓力

        ↕

👶 幼兒照顧（0-2歲）
需要托育、親職教育"""
    for para in left_frame.paragraphs:
        para.font.size = Pt(18)
        para.space_before = Pt(8)
        para.alignment = PP_ALIGN.CENTER

    # 右側：數據視覺化（大數字呈現）
    right_x = MARGIN_LEFT + Inches(6.0)

    # 數據1
    data1_num = slide.shapes.add_textbox(
        right_x, Inches(2.0), Inches(5.0), Inches(0.8)
    )
    data1_num.text_frame.text = "54.5%"
    data1_num.text_frame.paragraphs[0].font.size = Pt(72)
    data1_num.text_frame.paragraphs[0].font.bold = True
    data1_num.text_frame.paragraphs[0].font.color.rgb = MTK_ORANGE

    data1_desc = slide.shapes.add_textbox(
        right_x, Inches(2.9), Inches(5.0), Inches(0.4)
    )
    data1_desc.text_frame.text = "員工表示「時間貧窮」影響工作表現"
    data1_desc.text_frame.paragraphs[0].font.size = Pt(20)
    data1_desc.text_frame.paragraphs[0].font.color.rgb = DARK_GRAY

    # 數據2
    data2_num = slide.shapes.add_textbox(
        right_x, Inches(3.8), Inches(5.0), Inches(0.7)
    )
    data2_num.text_frame.text = "30 萬元"
    data2_num.text_frame.paragraphs[0].font.size = Pt(60)
    data2_num.text_frame.paragraphs[0].font.bold = True
    data2_num.text_frame.paragraphs[0].font.color.rgb = WARNING_ORANGE

    data2_desc = slide.shapes.add_textbox(
        right_x, Inches(4.6), Inches(5.0), Inches(0.4)
    )
    data2_desc.text_frame.text = "每位工程師離職成本（招募+訓練）"
    data2_desc.text_frame.paragraphs[0].font.size = Pt(20)
    data2_desc.text_frame.paragraphs[0].font.color.rgb = DARK_GRAY

    # 流程鏈
    flow_box = slide.shapes.add_textbox(
        right_x, Inches(5.4), Inches(5.0), Inches(1.0)
    )
    flow_frame = flow_box.text_frame
    flow_frame.text = "家庭照顧壓力  →  工作分心\n→  績效下降  →  選擇離職"
    for para in flow_frame.paragraphs:
        para.font.size = Pt(18)
        para.font.color.rgb = DARK_GRAY
        para.space_before = Pt(6)

    # ========== Slide 5: 挑戰2深入 - 招募困難 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(
        MARGIN_LEFT, MARGIN_TOP, CONTENT_WIDTH, Inches(0.7)
    )
    title_frame = title_box.text_frame
    title_frame.text = "挑戰 2：招募吸引力不足 → 人才外流"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = NAVY_BLUE

    # 對比表（竹科 vs 矽谷）
    compare_data = [
        ("薪資水平", "✓ 勝出", "❌ 略低"),
        ("生活成本", "✓ 較低", "❌ 高"),
        ("家庭支持", "❌ 缺乏", "✓ 完善"),
        ("工作文化", "⚠ 加班", "✓ 彈性"),
        ("職涯發展", "✓ 穩定", "✓ 多元")
    ]

    # 表格標題
    table_title_竹科 = slide.shapes.add_textbox(
        MARGIN_LEFT + Inches(2.5), Inches(1.8), Inches(3.0), Inches(0.5)
    )
    table_title_竹科.text_frame.text = "竹科（聯發科）"
    table_title_竹科.text_frame.paragraphs[0].font.size = Pt(24)
    table_title_竹科.text_frame.paragraphs[0].font.bold = True
    table_title_竹科.text_frame.paragraphs[0].font.color.rgb = MTK_ORANGE
    table_title_竹科.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    table_title_矽谷 = slide.shapes.add_textbox(
        MARGIN_LEFT + Inches(6.0), Inches(1.8), Inches(3.0), Inches(0.5)
    )
    table_title_矽谷.text_frame.text = "矽谷（三星/Intel）"
    table_title_矽谷.text_frame.paragraphs[0].font.size = Pt(24)
    table_title_矽谷.text_frame.paragraphs[0].font.bold = True
    table_title_矽谷.text_frame.paragraphs[0].font.color.rgb = NAVY_BLUE
    table_title_矽谷.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 對比行
    y_start = Inches(2.5)
    row_height = Inches(0.7)

    for i, (item, tsmc_status, competitor_status) in enumerate(compare_data):
        y_pos = y_start + i * row_height

        # 項目名稱
        item_tf = slide.shapes.add_textbox(
            MARGIN_LEFT, y_pos, Inches(2.0), Inches(0.5)
        )
        item_tf.text_frame.text = item
        item_tf.text_frame.paragraphs[0].font.size = Pt(20)
        item_tf.text_frame.paragraphs[0].font.bold = True
        item_tf.text_frame.paragraphs[0].font.color.rgb = DARK_GRAY

        # 竹科狀態
        tsmc_tf = slide.shapes.add_textbox(
            MARGIN_LEFT + Inches(2.5), y_pos, Inches(3.0), Inches(0.5)
        )
        tsmc_tf.text_frame.text = tsmc_status
        tsmc_tf.text_frame.paragraphs[0].font.size = Pt(22)
        tsmc_tf.text_frame.paragraphs[0].font.color.rgb = DARK_GRAY
        tsmc_tf.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # 競爭對手狀態
        comp_tf = slide.shapes.add_textbox(
            MARGIN_LEFT + Inches(6.0), y_pos, Inches(3.0), Inches(0.5)
        )
        comp_tf.text_frame.text = competitor_status
        comp_tf.text_frame.paragraphs[0].font.size = Pt(22)
        comp_tf.text_frame.paragraphs[0].font.color.rgb = DARK_GRAY
        comp_tf.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 底部關鍵發現
    finding_box = slide.shapes.add_textbox(
        MARGIN_LEFT, Inches(6.0), CONTENT_WIDTH, Inches(0.7)
    )
    finding_frame = finding_box.text_frame
    finding_frame.text = "⚠️  關鍵發現：「家庭支持」成為人才決策的關鍵因素\n年輕世代更重視 Work-Life Balance，不只看薪資"
    for para in finding_frame.paragraphs:
        para.font.size = Pt(20)
        para.font.bold = True
        para.font.color.rgb = MTK_ORANGE
        para.alignment = PP_ALIGN.CENTER
        para.space_before = Pt(4)

    # ========== Slide 6: 挑戰3深入 - ESG評級 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(
        MARGIN_LEFT, MARGIN_TOP, CONTENT_WIDTH, Inches(0.7)
    )
    title_frame = title_box.text_frame
    title_frame.text = "挑戰 3：ESG 評級落後 → 融資成本增加"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = NAVY_BLUE

    # ESG評級對比（視覺化）
    # 聯發科當前
    current_box = slide.shapes.add_textbox(
        MARGIN_LEFT + Inches(1.0), Inches(2.2), Inches(4.0), Inches(3.0)
    )
    current_frame = current_box.text_frame
    current_frame.text = "聯發科當前\n\nMSCI ESG 評級"
    current_para1 = current_frame.paragraphs[0]
    current_para1.font.size = Pt(22)
    current_para1.font.bold = True
    current_para1.font.color.rgb = DARK_GRAY
    current_para1.alignment = PP_ALIGN.CENTER

    # BBB大字
    p = current_frame.add_paragraph()
    p.text = "\nBBB"
    p.font.size = Pt(80)
    p.font.bold = True
    p.font.color.rgb = WARNING_ORANGE
    p.alignment = PP_ALIGN.CENTER

    p2 = current_frame.add_paragraph()
    p2.text = "\n產業平均: A"
    p2.font.size = Pt(18)
    p2.font.color.rgb = DARK_GRAY
    p2.alignment = PP_ALIGN.CENTER

    # 箭頭
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        MARGIN_LEFT + Inches(5.5), Inches(3.5), Inches(1.5), Inches(0.8)
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = SUCCESS_GREEN
    arrow.line.fill.background()

    # 目標
    target_box = slide.shapes.add_textbox(
        MARGIN_LEFT + Inches(7.5), Inches(2.2), Inches(4.0), Inches(3.0)
    )
    target_frame = target_box.text_frame
    target_frame.text = "改善目標\n\nMSCI ESG 評級"
    target_para1 = target_frame.paragraphs[0]
    target_para1.font.size = Pt(22)
    target_para1.font.bold = True
    target_para1.font.color.rgb = DARK_GRAY
    target_para1.alignment = PP_ALIGN.CENTER

    # A大字
    p = target_frame.add_paragraph()
    p.text = "\nA"
    p.font.size = Pt(96)
    p.font.bold = True
    p.font.color.rgb = SUCCESS_GREEN
    p.alignment = PP_ALIGN.CENTER

    p2 = target_frame.add_paragraph()
    p2.text = "\n提升 1 級"
    p2.font.size = Pt(18)
    p2.font.color.rgb = DARK_GRAY
    p2.alignment = PP_ALIGN.CENTER

    # 底部效益說明
    benefit_box = slide.shapes.add_textbox(
        MARGIN_LEFT, Inches(5.7), CONTENT_WIDTH, Inches(1.0)
    )
    benefit_frame = benefit_box.text_frame
    benefit_frame.text = "💰 效益計算\n\nESG 評級每提升 1 級 → 融資成本降低 0.5%\n聯發科年度融資規模 3,000 億 → 每年節省 15 億元"
    for i, para in enumerate(benefit_frame.paragraphs):
        if i == 0:
            para.font.size = Pt(20)
            para.font.bold = True
            para.font.color.rgb = SUCCESS_GREEN
        else:
            para.font.size = Pt(18)
            para.font.color.rgb = DARK_GRAY
        para.alignment = PP_ALIGN.CENTER
        para.space_before = Pt(4)

    # ========== Slide 7: 解決方案 - 照護時數銀行 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(
        MARGIN_LEFT, MARGIN_TOP, CONTENT_WIDTH, Inches(0.7)
    )
    title_frame = title_box.text_frame
    title_frame.text = "解決方案：竹科家庭照護時數銀行"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = NAVY_BLUE

    # 循環流程圖（中央圓形）
    center_x = MARGIN_LEFT + Inches(4.5)
    center_y = Inches(3.5)

    # 中心圓：赤土崎樞紐
    center_circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        center_x, center_y, Inches(2.5), Inches(2.5)
    )
    center_circle.fill.solid()
    center_circle.fill.fore_color.rgb = SUCCESS_GREEN
    center_circle.line.color.rgb = NAVY_BLUE
    center_circle.line.width = Pt(3)
    center_tf = center_circle.text_frame
    center_tf.text = "赤土崎\n全齡樞紐\n\n8,000\n小時/年"
    center_tf.paragraphs[0].font.size = Pt(24)
    center_tf.paragraphs[0].font.bold = True
    center_tf.paragraphs[0].font.color.rgb = WHITE
    center_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    center_tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 步驟1：聯發科投資（左上）
    step1 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        MARGIN_LEFT, Inches(2.0), Inches(3.5), Inches(1.2)
    )
    step1.fill.solid()
    step1.fill.fore_color.rgb = MTK_ORANGE
    step1.line.fill.background()
    step1_tf = step1.text_frame
    step1_tf.text = "① 聯發科投資\n\n400萬元/年"
    for para in step1_tf.paragraphs:
        para.font.size = Pt(20) if para.text.startswith("①") else Pt(28)
        para.font.bold = True
        para.font.color.rgb = WHITE
        para.alignment = PP_ALIGN.CENTER
    step1_tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 步驟3：員工兌換（右上）
    step3 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        MARGIN_LEFT + Inches(7.8), Inches(2.0), Inches(3.5), Inches(1.2)
    )
    step3.fill.solid()
    step3.fill.fore_color.rgb = NAVY_BLUE
    step3.line.fill.background()
    step3_tf = step3.text_frame
    step3_tf.text = "② 員工憑證兌換\n\n1,700人次/年"
    for para in step3_tf.paragraphs:
        para.font.size = Pt(20) if para.text.startswith("②") else Pt(26)
        para.font.bold = True
        para.font.color.rgb = WHITE
        para.alignment = PP_ALIGN.CENTER
    step3_tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 步驟4：三重回報（底部）
    step4 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        MARGIN_LEFT + Inches(2.5), Inches(5.8), Inches(6.3), Inches(0.9)
    )
    step4.fill.solid()
    step4.fill.fore_color.rgb = SUCCESS_GREEN
    step4.line.fill.background()
    step4_tf = step4.text_frame
    step4_tf.text = "③ 聯發科獲得：留任率+5% | ESG評級↑ | SROI 1:2.04"
    step4_tf.paragraphs[0].font.size = Pt(22)
    step4_tf.paragraphs[0].font.bold = True
    step4_tf.paragraphs[0].font.color.rgb = WHITE
    step4_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    step4_tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # ========== Slide 8: SROI 大數字 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(
        MARGIN_LEFT, MARGIN_TOP, CONTENT_WIDTH, Inches(0.7)
    )
    title_frame = title_box.text_frame
    title_frame.text = "投資效益：每投入 1 元，創造 2.04 元社會價值"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(38)
    title_para.font.bold = True
    title_para.font.color.rgb = NAVY_BLUE

    # 視覺化流程
    # 投資400萬
    invest_box = slide.shapes.add_textbox(
        MARGIN_LEFT + Inches(1.0), Inches(2.5), Inches(3.0), Inches(1.5)
    )
    invest_frame = invest_box.text_frame
    invest_frame.text = "投資\n\n500 萬元"
    invest_frame.paragraphs[0].font.size = Pt(24)
    invest_frame.paragraphs[0].font.color.rgb = DARK_GRAY
    invest_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = invest_frame.paragraphs[1]
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = MTK_ORANGE
    p.alignment = PP_ALIGN.CENTER
    invest_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 乘號
    multiply = slide.shapes.add_textbox(
        MARGIN_LEFT + Inches(4.2), Inches(2.8), Inches(0.8), Inches(1.0)
    )
    multiply.text_frame.text = "×"
    multiply.text_frame.paragraphs[0].font.size = Pt(72)
    multiply.text_frame.paragraphs[0].font.bold = True
    multiply.text_frame.paragraphs[0].font.color.rgb = SUCCESS_GREEN
    multiply.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    multiply.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # SROI倍數
    sroi_box = slide.shapes.add_textbox(
        MARGIN_LEFT + Inches(5.2), Inches(2.5), Inches(2.5), Inches(1.5)
    )
    sroi_frame = sroi_box.text_frame
    sroi_frame.text = "SROI\n\n2.04"
    sroi_frame.paragraphs[0].font.size = Pt(24)
    sroi_frame.paragraphs[0].font.color.rgb = DARK_GRAY
    sroi_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = sroi_frame.paragraphs[1]
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = SUCCESS_GREEN
    p.alignment = PP_ALIGN.CENTER
    sroi_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 等號
    equals = slide.shapes.add_textbox(
        MARGIN_LEFT + Inches(7.9), Inches(2.8), Inches(0.8), Inches(1.0)
    )
    equals.text_frame.text = "="
    equals.text_frame.paragraphs[0].font.size = Pt(72)
    equals.text_frame.paragraphs[0].font.bold = True
    equals.text_frame.paragraphs[0].font.color.rgb = DARK_GRAY
    equals.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    equals.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 社會價值
    value_box = slide.shapes.add_textbox(
        MARGIN_LEFT + Inches(8.9), Inches(2.5), Inches(3.0), Inches(1.5)
    )
    value_frame = value_box.text_frame
    value_frame.text = "社會價值\n\n1,020 萬"
    value_frame.paragraphs[0].font.size = Pt(24)
    value_frame.paragraphs[0].font.color.rgb = DARK_GRAY
    value_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = value_frame.paragraphs[1]
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = SUCCESS_GREEN
    p.alignment = PP_ALIGN.CENTER
    value_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 底部說明
    note_box = slide.shapes.add_textbox(
        MARGIN_LEFT, Inches(5.0), CONTENT_WIDTH, Inches(1.2)
    )
    note_frame = note_box.text_frame
    note_frame.text = "📊 計算標準\n\n依據「台灣社會影響力研究院」SROI 認證標準\n評估期間：5年期，折現率 3.5%"
    note_frame.paragraphs[0].font.size = Pt(20)
    note_frame.paragraphs[0].font.bold = True
    note_frame.paragraphs[0].font.color.rgb = NAVY_BLUE
    note_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    for para in note_frame.paragraphs[1:]:
        para.font.size = Pt(16)
        para.font.color.rgb = DARK_GRAY
        para.alignment = PP_ALIGN.CENTER
        para.space_before = Pt(6)

    # ========== Slide 9: 建築規劃（視覺化剖面圖）==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(
        MARGIN_LEFT, MARGIN_TOP, CONTENT_WIDTH, Inches(0.7)
    )
    title_frame = title_box.text_frame
    title_frame.text = "赤土崎全齡社福樞紐 - 一站式整合服務"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = NAVY_BLUE

    # 建築剖面圖（分層視覺化）
    floors = [
        ("4F", "青少年 STEAM 中心", "500m²", "120人次/月", LIGHT_BLUE),
        ("3F", "家庭支持中心", "500m²", "80家庭/月", RGBColor(255, 213, 79)),
        ("2F", "公共托嬰中心", "700m²", "45人（聯發科優先）", RGBColor(255, 183, 77)),
        ("1F", "長照日照中心", "800m²", "55人（聯發科優先）", RGBColor(149, 117, 205)),
        ("B1", "停車場+設備", "600m²", "30車位（聯發科保留10）", DARK_GRAY),
    ]

    x_start = MARGIN_LEFT + Inches(1.5)
    y_start = Inches(1.8)
    floor_height = Inches(0.85)
    floor_width = Inches(8.5)

    for i, (floor_num, floor_name, area, capacity, color) in enumerate(floors):
        y_pos = y_start + i * floor_height

        # 樓層色塊
        floor_rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x_start, y_pos, floor_width, floor_height
        )
        floor_rect.fill.solid()
        floor_rect.fill.fore_color.rgb = color
        floor_rect.line.color.rgb = WHITE
        floor_rect.line.width = Pt(2)

        # 樓層編號
        floor_num_tf = slide.shapes.add_textbox(
            x_start + Inches(0.2), y_pos + Inches(0.15),
            Inches(0.6), Inches(0.5)
        )
        floor_num_tf.text_frame.text = floor_num
        floor_num_tf.text_frame.paragraphs[0].font.size = Pt(24)
        floor_num_tf.text_frame.paragraphs[0].font.bold = True
        floor_num_tf.text_frame.paragraphs[0].font.color.rgb = WHITE if floor_num != "B1" else LIGHT_GRAY
        floor_num_tf.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        # 樓層名稱
        floor_name_tf = slide.shapes.add_textbox(
            x_start + Inches(1.0), y_pos + Inches(0.15),
            Inches(3.5), Inches(0.5)
        )
        floor_name_tf.text_frame.text = floor_name
        floor_name_tf.text_frame.paragraphs[0].font.size = Pt(20)
        floor_name_tf.text_frame.paragraphs[0].font.bold = True
        floor_name_tf.text_frame.paragraphs[0].font.color.rgb = WHITE if floor_num != "B1" else LIGHT_GRAY
        floor_name_tf.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        # 面積
        area_tf = slide.shapes.add_textbox(
            x_start + Inches(4.7), y_pos + Inches(0.15),
            Inches(1.2), Inches(0.5)
        )
        area_tf.text_frame.text = area
        area_tf.text_frame.paragraphs[0].font.size = Pt(16)
        area_tf.text_frame.paragraphs[0].font.color.rgb = WHITE if floor_num != "B1" else LIGHT_GRAY
        area_tf.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        # 容量
        capacity_tf = slide.shapes.add_textbox(
            x_start + Inches(6.0), y_pos + Inches(0.15),
            Inches(2.3), Inches(0.5)
        )
        capacity_tf.text_frame.text = capacity
        capacity_tf.text_frame.paragraphs[0].font.size = Pt(16)
        capacity_tf.text_frame.paragraphs[0].font.color.rgb = WHITE if floor_num != "B1" else LIGHT_GRAY
        capacity_tf.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 總計標註
    total_box = slide.shapes.add_textbox(
        MARGIN_LEFT, Inches(6.3), CONTENT_WIDTH, Inches(0.5)
    )
    total_frame = total_box.text_frame
    total_frame.text = "總建築面積：3,100m² | 地上4層+地下1層 | 距離竹科 2.5km（車程5分鐘）"
    total_frame.paragraphs[0].font.size = Pt(18)
    total_frame.paragraphs[0].font.bold = True
    total_frame.paragraphs[0].font.color.rgb = NAVY_BLUE
    total_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # ========== Slide 10: 投資方案 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(
        MARGIN_LEFT, MARGIN_TOP, CONTENT_WIDTH, Inches(0.7)
    )
    title_frame = title_box.text_frame
    title_frame.text = "聯發科投資方案 - 三年期合作計畫"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = NAVY_BLUE

    # 時間軸視覺化
    years = ["2025", "2026", "2027"]
    amounts = ["400萬", "400萬", "400萬"]
    phases = ["設計階段", "建設階段", "營運階段"]

    timeline_y = Inches(2.5)
    card_width = Inches(3.3)
    card_height = Inches(2.8)
    spacing = Inches(0.5)

    for i, (year, amount, phase) in enumerate(zip(years, amounts, phases)):
        x_pos = MARGIN_LEFT + Inches(0.5) + i * (card_width + spacing)

        # 卡片
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x_pos, timeline_y, card_width, card_height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_GRAY
        card.line.color.rgb = MTK_ORANGE
        card.line.width = Pt(2)

        # 年份
        year_tf = slide.shapes.add_textbox(
            x_pos + Inches(0.2), timeline_y + Inches(0.2),
            card_width - Inches(0.4), Inches(0.5)
        )
        year_tf.text_frame.text = year
        year_tf.text_frame.paragraphs[0].font.size = Pt(32)
        year_tf.text_frame.paragraphs[0].font.bold = True
        year_tf.text_frame.paragraphs[0].font.color.rgb = MTK_ORANGE
        year_tf.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # 金額
        amount_tf = slide.shapes.add_textbox(
            x_pos + Inches(0.2), timeline_y + Inches(0.9),
            card_width - Inches(0.4), Inches(0.7)
        )
        amount_tf.text_frame.text = amount
        amount_tf.text_frame.paragraphs[0].font.size = Pt(48)
        amount_tf.text_frame.paragraphs[0].font.bold = True
        amount_tf.text_frame.paragraphs[0].font.color.rgb = NAVY_BLUE
        amount_tf.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # 階段
        phase_tf = slide.shapes.add_textbox(
            x_pos + Inches(0.2), timeline_y + Inches(1.8),
            card_width - Inches(0.4), Inches(0.8)
        )
        phase_tf.text_frame.text = phase
        phase_tf.text_frame.paragraphs[0].font.size = Pt(20)
        phase_tf.text_frame.paragraphs[0].font.color.rgb = DARK_GRAY
        phase_tf.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        phase_tf.text_frame.word_wrap = True

    # 總計
    total_box = slide.shapes.add_textbox(
        MARGIN_LEFT + Inches(2.5), Inches(5.6), Inches(6.3), Inches(0.8)
    )
    total_frame = total_box.text_frame
    total_frame.text = "3 年總投資：1,500 萬元"
    total_frame.paragraphs[0].font.size = Pt(40)
    total_frame.paragraphs[0].font.bold = True
    total_frame.paragraphs[0].font.color.rgb = SUCCESS_GREEN
    total_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    total_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # ========== Slide 11: 行動呼籲 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 背景色
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_GRAY

    # 標題
    title_box = slide.shapes.add_textbox(
        MARGIN_LEFT, MARGIN_TOP, CONTENT_WIDTH, Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = "我們需要聯發科的決定"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = NAVY_BLUE
    title_para.alignment = PP_ALIGN.CENTER

    # 中央行動呼籲框
    cta_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        MARGIN_LEFT + Inches(1.5), Inches(2.2), Inches(8.3), Inches(2.5)
    )
    cta_box.fill.solid()
    cta_box.fill.fore_color.rgb = MTK_ORANGE
    cta_box.line.fill.background()

    cta_text = slide.shapes.add_textbox(
        MARGIN_LEFT + Inches(2.0), Inches(2.5), Inches(7.3), Inches(2.0)
    )
    cta_frame = cta_text.text_frame
    cta_frame.text = "請在 11/15 前決定\n\n是否參與「竹科家庭照護時數銀行」\n投資額度：400萬/年 × 3年"
    cta_frame.paragraphs[0].font.size = Pt(32)
    cta_frame.paragraphs[0].font.bold = True
    cta_frame.paragraphs[0].font.color.rgb = WHITE
    cta_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    for para in cta_frame.paragraphs[1:]:
        para.font.size = Pt(24)
        para.font.color.rgb = WHITE
        para.alignment = PP_ALIGN.CENTER
        para.space_before = Pt(10)
    cta_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 時機提醒
    urgency_box = slide.shapes.add_textbox(
        MARGIN_LEFT, Inches(5.2), CONTENT_WIDTH, Inches(0.8)
    )
    urgency_frame = urgency_box.text_frame
    urgency_frame.text = "⏰ 黃金窗口：10/30 標案開標後，11/8-11/28 設計階段\n這是納入「聯發科專區」的最後機會"
    urgency_frame.paragraphs[0].font.size = Pt(20)
    urgency_frame.paragraphs[0].font.bold = True
    urgency_frame.paragraphs[0].font.color.rgb = WARNING_ORANGE
    urgency_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    for para in urgency_frame.paragraphs[1:]:
        para.font.size = Pt(18)
        para.font.color.rgb = DARK_GRAY
        para.alignment = PP_ALIGN.CENTER
        para.space_before = Pt(4)

    # 聯絡資訊
    contact_box = slide.shapes.add_textbox(
        MARGIN_LEFT, Inches(6.4), CONTENT_WIDTH, Inches(0.6)
    )
    contact_frame = contact_box.text_frame
    contact_frame.text = "📞 專案顧問：Easterlin Hsinchu Team\n🏛️ 新竹市政府社會處"
    for para in contact_frame.paragraphs:
        para.font.size = Pt(18)
        para.font.color.rgb = DARK_GRAY
        para.alignment = PP_ALIGN.CENTER
        para.space_before = Pt(4)

    # 保存簡報
    filename = "聯發科ESG合作提案_V2_16x9_2025.pptx"
    prs.save(filename)
    print(f"[SUCCESS] PowerPoint V2.0 created: {filename}")
    print(f"[INFO] Slide format: 16:9 (13.333\" × 7.5\")")
    print(f"[INFO] Design principle: Visual-first, One-slide-one-concept")
    print(f"[INFO] Total slides: 11")
    print(f"[INFO] Key improvements:")
    print(f"       - 16:9 aspect ratio ✓")
    print(f"       - Visual hierarchy ✓")
    print(f"       - Large fonts (24-96pt) ✓")
    print(f"       - Data visualization ✓")
    print(f"       - Reduced text density ✓")

    return filename

if __name__ == "__main__":
    create_mediatek_presentation_v2()
